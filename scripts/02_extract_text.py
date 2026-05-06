from __future__ import annotations

import json
from pathlib import Path
from typing import Any

import fitz
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation

from rag_common import (
    ensure_runtime_dirs,
    jsonl_read,
    jsonl_write,
    load_config,
    normalize_text,
    read_text_guess,
    resolve_work_path,
    safe_rel_id,
)


TEXT_EXTENSIONS = {".md", ".txt", ".json", ".yml", ".yaml", ".drawio", ".puml", ".srt", ".py", ".js", ".ts", ".css"}


def extract_docx(path: Path) -> list[dict[str, Any]]:
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table_index, table in enumerate(doc.tables, start=1):
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        if rows:
            parts.append(f"\n[Table {table_index}]\n" + "\n".join(rows))
    return [{"section": "document", "text": "\n".join(parts)}]


def extract_pdf(path: Path) -> list[dict[str, Any]]:
    rows = []
    with fitz.open(str(path)) as doc:
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text.strip():
                rows.append({"section": f"page {page_index}", "page": page_index, "text": text})
    return rows


def extract_pptx(path: Path) -> list[dict[str, Any]]:
    prs = Presentation(str(path))
    rows = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
        if parts:
            rows.append({"section": f"slide {slide_index}", "slide": slide_index, "text": "\n".join(parts)})
    return rows


def extract_excel(path: Path) -> list[dict[str, Any]]:
    rows = []
    engine = "pyxlsb" if path.suffix.lower() == ".xlsb" else None
    sheets = pd.read_excel(path, sheet_name=None, dtype=str, engine=engine)
    for sheet_name, df in sheets.items():
        df = df.fillna("")
        text = df.to_csv(index=False, sep="\t")
        if text.strip():
            rows.append({"section": f"sheet {sheet_name}", "sheet": str(sheet_name), "text": text})
    return rows


def extract_html(path: Path) -> list[dict[str, Any]]:
    soup = BeautifulSoup(read_text_guess(path), "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text("\n")
    return [{"section": "html", "text": text}]


def extract_text_file(path: Path) -> list[dict[str, Any]]:
    return [{"section": "file", "text": read_text_guess(path)}]


def extract_one(path: Path) -> list[dict[str, Any]]:
    ext = path.suffix.lower()
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext in {".xlsx", ".xlsb"}:
        return extract_excel(path)
    if ext == ".html":
        return extract_html(path)
    if ext in TEXT_EXTENSIONS:
        return extract_text_file(path)
    return []


def main() -> None:
    cfg = load_config()
    ensure_runtime_dirs(cfg)

    manifest_path = resolve_work_path(cfg, cfg["paths"]["manifest"])
    extracted_dir = resolve_work_path(cfg, cfg["paths"]["extracted_text_dir"])
    metadata_path = extracted_dir / "_metadata.jsonl"

    metadata = []
    ok = 0
    errors = 0
    skipped = 0

    for rec in jsonl_read(manifest_path):
        if rec.get("status") != "included":
            skipped += 1
            continue

        source = Path(rec["path"])
        out_id = safe_rel_id(rec["relative_path"])
        out_path = extracted_dir / f"{out_id}.{rec['sha256'][:12]}.txt"

        try:
            sections = extract_one(source)
            text_parts = []
            for section in sections:
                text = normalize_text(section.get("text", ""))
                if text:
                    text_parts.append(f"# {section.get('section', 'section')}\n\n{text}")

            text = "\n\n".join(text_parts).strip()
            if not text:
                skipped += 1
                continue

            out_path.write_text(text, encoding="utf-8", newline="\n")
            metadata.append(
                {
                    "source_path": rec["path"],
                    "relative_path": rec["relative_path"],
                    "extension": rec["extension"],
                    "sha256": rec["sha256"],
                    "mtime": rec["mtime"],
                    "extracted_path": str(out_path),
                    "chars": len(text),
                }
            )
            ok += 1
        except Exception as exc:
            errors += 1
            metadata.append(
                {
                    "source_path": rec["path"],
                    "relative_path": rec["relative_path"],
                    "extension": rec["extension"],
                    "sha256": rec.get("sha256"),
                    "error": repr(exc),
                }
            )

    jsonl_write(metadata_path, metadata)
    print(json.dumps({"extracted": ok, "skipped": skipped, "errors": errors, "metadata": str(metadata_path)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
