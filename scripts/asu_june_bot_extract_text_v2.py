from __future__ import annotations

import argparse
import json
import shutil
import sys
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable

import fitz
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation


WORK_ROOT = Path(__file__).resolve().parents[1]
SRC_DIR = WORK_ROOT / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

from asu_june_bot.core.config import load_config, resolve_work_path  # noqa: E402
from asu_june_bot.ingestion.models import ExtractedBlock, SourceDocument  # noqa: E402
from asu_june_bot.ingestion.utils import (  # noqa: E402
    detect_heading_level,
    enrich_block_metadata,
    iter_source_files,
    jsonl_write,
    normalize_text,
    read_text_guess,
    relative_to_project,
    sha256_file,
    stable_id,
)


EXTRACTOR_VERSION = "v2.1"
DEFAULT_OUTPUT_DIR = "data/asu_june_bot/extracted_v2"
TEXT_EXTENSIONS = {".md", ".txt", ".json", ".yml", ".yaml", ".drawio", ".puml", ".bpmn", ".srt"}
MAX_EXCEL_COLUMNS = 120


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def append_jsonl(path: Path, rows: Iterable[dict[str, Any]]) -> int:
    path.parent.mkdir(parents=True, exist_ok=True)
    count = 0
    with path.open("a", encoding="utf-8", newline="\n") as fp:
        for row in rows:
            fp.write(json.dumps(row, ensure_ascii=False) + "\n")
            count += 1
        fp.flush()
    return count


def read_jsonl_if_exists(path: Path) -> list[dict[str, Any]]:
    if not path.exists():
        return []
    rows: list[dict[str, Any]] = []
    with path.open("r", encoding="utf-8") as fp:
        for line in fp:
            line = line.strip()
            if not line:
                continue
            try:
                rows.append(json.loads(line))
            except json.JSONDecodeError:
                # Keep resume robust against a truncated last line.
                continue
    return rows


def load_completed_source_ids(documents_path: Path) -> set[str]:
    rows = read_jsonl_if_exists(documents_path)
    return {str(row.get("source_id")) for row in rows if row.get("source_id")}


def load_source_links(path: Path) -> dict[str, str]:
    links: dict[str, str] = {}
    for row in read_jsonl_if_exists(path):
        rel = str(row.get("relative_path") or "").replace("\\", "/").strip("/")
        url = str(row.get("source_url") or "").strip()
        if rel and url:
            links[rel] = url
    return links


def make_source_document(cfg: dict[str, Any], path: Path) -> SourceDocument:
    relative_path = relative_to_project(cfg, path)
    sha = sha256_file(path)
    source_id = stable_id(f"source:{relative_path}:{sha}", length=32)
    probe_meta = enrich_block_metadata(relative_path, path.suffix.lower(), relative_path)
    source_links = cfg.get("source_links") if isinstance(cfg.get("source_links"), dict) else {}
    source_url = source_links.get(relative_path.replace("\\", "/")) if source_links else None
    stat = path.stat()
    return SourceDocument(
        source_id=source_id,
        source_path=str(path.resolve()),
        relative_path=relative_path,
        extension=path.suffix.lower(),
        sha256=sha,
        mtime=stat.st_mtime,
        size_bytes=stat.st_size,
        source_type=probe_meta.get("source_type"),
        document_type=probe_meta.get("document_type"),
        stage=probe_meta.get("stage"),
        module=probe_meta.get("module"),
        source_url=source_url,
    )


def make_block(
    *,
    source: SourceDocument,
    block_index: int,
    block_type: str,
    text: str,
    page: int | None = None,
    slide: int | None = None,
    sheet: str | None = None,
    paragraph_index: int | None = None,
    heading_level: int | None = None,
    style_name: str | None = None,
    table_index: int | None = None,
    row_index: int | None = None,
    headers: list[str] | None = None,
    cells: dict[str, str] | None = None,
    title: str | None = None,
    parent_hint: str | None = None,
) -> ExtractedBlock | None:
    clean = normalize_text(text)
    if not clean:
        return None
    enriched = enrich_block_metadata(source.relative_path, source.extension, clean)
    table_id = f"Table {table_index}" if table_index is not None else None
    row_id = str(row_index) if row_index is not None else None
    block_id = stable_id(
        f"block:{source.sha256}:{block_index}:{block_type}:{page or ''}:{slide or ''}:{sheet or ''}:{table_id or ''}:{row_id or ''}:{clean[:160]}",
        length=32,
    )
    return ExtractedBlock(
        block_id=block_id,
        source_id=source.source_id,
        block_index=block_index,
        block_type=block_type,
        text=clean,
        source_path=source.source_path,
        relative_path=source.relative_path,
        extension=source.extension,
        sha256=source.sha256,
        mtime=source.mtime,
        document_name=Path(source.relative_path).name,
        source_type=enriched.get("source_type"),
        document_type=enriched.get("document_type"),
        stage=enriched.get("stage"),
        module=enriched.get("module"),
        source_url=source.source_url,
        page=page,
        slide=slide,
        sheet=sheet,
        paragraph_index=paragraph_index,
        heading_level=heading_level,
        style_name=style_name,
        section=enriched.get("section"),
        sections=enriched.get("sections") or [],
        table_id=table_id,
        table_index=table_index,
        row_id=row_id,
        row_index=row_index,
        col_count=len(headers or []) if headers else (len(cells or {}) if cells else None),
        headers=headers or [],
        cells=cells or {},
        title=title,
        parent_hint=parent_hint,
    )


def iter_docx_blocks(document: Document) -> Iterable[Paragraph | Table]:
    body = document.element.body
    for child in body.iterchildren():
        if child.tag.endswith("}p"):
            yield Paragraph(child, document)
        elif child.tag.endswith("}tbl"):
            yield Table(child, document)


def table_row_values(table: Table, row_index: int) -> list[str]:
    row = table.rows[row_index]
    return [normalize_text(cell.text) for cell in row.cells]


def unique_headers(values: list[str]) -> list[str]:
    headers: list[str] = []
    seen: Counter[str] = Counter()
    for idx, value in enumerate(values, start=1):
        header = value.strip() or f"col_{idx}"
        seen[header] += 1
        if seen[header] > 1:
            header = f"{header}_{seen[header]}"
        headers.append(header)
    return headers


def trim_empty_columns(rows: list[list[str]], max_cols: int = MAX_EXCEL_COLUMNS) -> list[list[str]]:
    """Keep only the useful part of a worksheet used range.

    Some Excel files report thousands of formatted columns as part of the used
    range. Those empty/generated columns later become col_N headers repeated in
    every table_row chunk, which can inflate one row into hundreds of thousands
    of characters.
    """
    if not rows:
        return []
    last_non_empty = -1
    for row in rows:
        for idx, value in enumerate(row):
            if str(value or "").strip():
                last_non_empty = max(last_non_empty, idx)
    if last_non_empty < 0:
        return []
    width = min(last_non_empty + 1, max_cols)
    return [row[:width] for row in rows]


def compact_table_cells(headers: list[str], values: list[str]) -> dict[str, str]:
    cells: dict[str, str] = {}
    for idx, value in enumerate(values[: len(headers)]):
        clean_value = normalize_text(value)
        if not clean_value:
            continue
        header = headers[idx] if idx < len(headers) else f"col_{idx + 1}"
        cells[header] = clean_value
    return cells


def row_text_from_cells(row_label: str, cells: dict[str, str]) -> str:
    if not cells:
        return ""
    pairs = [f"{key}: {value}" for key, value in cells.items()]
    return row_label + "\n" + "\n".join(pairs)


def header_score(values: list[str]) -> int:
    text = " ".join(values).lower()
    keywords = (
        "№",
        "номер",
        "код",
        "наименование",
        "описание",
        "требование",
        "атрибут",
        "тип",
        "значение",
        "отправитель",
        "получатель",
        "система",
        "статус",
        "дата",
        "роль",
    )
    score = sum(1 for value in values if normalize_text(value))
    score += sum(3 for keyword in keywords if keyword in text)
    return score


def detect_header_row(rows: list[list[str]], max_scan_rows: int = 8) -> int:
    if not rows:
        return 0
    scan = rows[:max_scan_rows]
    best_idx = 0
    best_score = -1
    for idx, values in enumerate(scan):
        score = header_score(values)
        if score > best_score:
            best_idx = idx
            best_score = score
    return best_idx


def extract_docx_v2(path: Path, source: SourceDocument) -> list[ExtractedBlock]:
    doc = Document(str(path))
    blocks: list[ExtractedBlock] = []
    block_index = 0
    paragraph_index = 0
    table_index = 0
    current_parent_hint: str | None = None

    for item in iter_docx_blocks(doc):
        if isinstance(item, Paragraph):
            paragraph_index += 1
            text = normalize_text(item.text)
            if not text:
                continue
            style_name = item.style.name if item.style is not None else None
            heading_level = detect_heading_level(style_name)
            if heading_level is not None:
                current_parent_hint = text
            block = make_block(
                source=source,
                block_index=block_index,
                block_type="heading" if heading_level is not None else "paragraph",
                text=text,
                paragraph_index=paragraph_index,
                heading_level=heading_level,
                style_name=style_name,
                title=text if heading_level is not None else None,
                parent_hint=current_parent_hint,
            )
            if block:
                blocks.append(block)
                block_index += 1
            continue

        if isinstance(item, Table):
            table_index += 1
            if not item.rows:
                continue
            table_rows = [table_row_values(item, idx) for idx in range(len(item.rows))]
            header_idx = detect_header_row(table_rows)
            headers = unique_headers(table_rows[header_idx])
            table_title = current_parent_hint
            table_text_lines = [f"Таблица {table_index}"]
            if table_title:
                table_text_lines.append(f"Контекст: {table_title}")
            table_text_lines.append("Заголовки: " + " | ".join(headers))

            parent_block = make_block(
                source=source,
                block_index=block_index,
                block_type="table",
                text="\n".join(table_text_lines),
                table_index=table_index,
                headers=headers,
                title=table_title,
                parent_hint=current_parent_hint,
            )
            if parent_block:
                blocks.append(parent_block)
                block_index += 1

            for idx, values in enumerate(table_rows):
                if idx == header_idx:
                    continue
                if not any(values):
                    continue
                cells = compact_table_cells(headers, values)
                if not cells:
                    continue
                row_text = "\n".join(
                    [
                        f"Таблица {table_index}",
                        f"Контекст: {table_title}" if table_title else "",
                        row_text_from_cells(f"Строка {idx + 1}:", cells),
                    ]
                )
                block = make_block(
                    source=source,
                    block_index=block_index,
                    block_type="table_row",
                    text=row_text,
                    table_index=table_index,
                    row_index=idx + 1,
                    headers=headers,
                    cells=cells,
                    title=table_title,
                    parent_hint=current_parent_hint,
                )
                if block:
                    blocks.append(block)
                    block_index += 1
    return blocks


def normalize_cell(value: Any) -> str:
    if value is None:
        return ""
    return normalize_text(str(value))


def extract_xlsx_with_openpyxl(path: Path, source: SourceDocument) -> list[ExtractedBlock]:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    blocks: list[ExtractedBlock] = []
    block_index = 0

    for sheet in workbook.worksheets:
        raw_rows = [[normalize_cell(cell) for cell in row] for row in sheet.iter_rows(values_only=True)]
        raw_rows = [row for row in raw_rows if any(row)]
        raw_rows = trim_empty_columns(raw_rows)
        raw_rows = [row for row in raw_rows if any(row)]
        if not raw_rows:
            continue
        header_idx = detect_header_row(raw_rows)
        headers = unique_headers(raw_rows[header_idx])
        if len(headers) >= MAX_EXCEL_COLUMNS:
            print(
                f"[extract:xlsx] {source.relative_path} / {sheet.title}: worksheet width capped at {MAX_EXCEL_COLUMNS} columns",
                file=sys.stderr,
            )
        sheet_name = str(sheet.title)
        sheet_text = f"Лист: {sheet_name}\nЗаголовки: " + " | ".join(headers)
        block = make_block(
            source=source,
            block_index=block_index,
            block_type="sheet",
            text=sheet_text,
            sheet=sheet_name,
            headers=headers,
            title=sheet_name,
        )
        if block:
            blocks.append(block)
            block_index += 1

        for row_idx, values in enumerate(raw_rows, start=1):
            if row_idx - 1 == header_idx:
                continue
            cells = compact_table_cells(headers, values)
            if not cells:
                continue
            row_text = "\n".join(
                [
                    f"Лист: {sheet_name}",
                    row_text_from_cells(f"Строка {row_idx}:", cells),
                ]
            )
            block = make_block(
                source=source,
                block_index=block_index,
                block_type="table_row",
                text=row_text,
                sheet=sheet_name,
                row_index=row_idx,
                headers=headers,
                cells=cells,
                title=sheet_name,
                parent_hint=sheet_name,
            )
            if block:
                blocks.append(block)
                block_index += 1
    workbook.close()
    return blocks


def extract_xlsb_with_pandas(path: Path, source: SourceDocument) -> list[ExtractedBlock]:
    sheets = pd.read_excel(path, sheet_name=None, dtype=str, engine="pyxlsb", header=None)
    blocks: list[ExtractedBlock] = []
    block_index = 0

    for sheet_name, df in sheets.items():
        df = df.fillna("")
        raw_rows = [[normalize_cell(value) for value in row] for row in df.values.tolist()]
        raw_rows = [row for row in raw_rows if any(row)]
        raw_rows = trim_empty_columns(raw_rows)
        raw_rows = [row for row in raw_rows if any(row)]
        if not raw_rows:
            continue
        header_idx = detect_header_row(raw_rows)
        headers = unique_headers(raw_rows[header_idx])
        if len(headers) >= MAX_EXCEL_COLUMNS:
            print(
                f"[extract:xlsb] {source.relative_path} / {sheet_name}: worksheet width capped at {MAX_EXCEL_COLUMNS} columns",
                file=sys.stderr,
            )
        sheet_title = str(sheet_name)
        sheet_text = f"Лист: {sheet_title}\nЗаголовки: " + " | ".join(headers)
        block = make_block(
            source=source,
            block_index=block_index,
            block_type="sheet",
            text=sheet_text,
            sheet=sheet_title,
            headers=headers,
            title=sheet_title,
        )
        if block:
            blocks.append(block)
            block_index += 1

        for row_idx, values in enumerate(raw_rows, start=1):
            if row_idx - 1 == header_idx:
                continue
            cells = compact_table_cells(headers, values)
            if not cells:
                continue
            row_text = "\n".join(
                [
                    f"Лист: {sheet_title}",
                    row_text_from_cells(f"Строка {row_idx}:", cells),
                ]
            )
            block = make_block(
                source=source,
                block_index=block_index,
                block_type="table_row",
                text=row_text,
                sheet=sheet_title,
                row_index=row_idx,
                headers=headers,
                cells=cells,
                title=sheet_title,
                parent_hint=sheet_title,
            )
            if block:
                blocks.append(block)
                block_index += 1
    return blocks


def extract_xlsx_v2(path: Path, source: SourceDocument) -> list[ExtractedBlock]:
    if path.suffix.lower() == ".xlsb":
        return extract_xlsb_with_pandas(path, source)
    return extract_xlsx_with_openpyxl(path, source)


def extract_pdf_v2(path: Path, source: SourceDocument) -> list[ExtractedBlock]:
    blocks: list[ExtractedBlock] = []
    block_index = 0
    with fitz.open(str(path)) as doc:
        for page_index, page in enumerate(doc, start=1):
            text = normalize_text(page.get_text("text"))
            block = make_block(
                source=source,
                block_index=block_index,
                block_type="page",
                text=text,
                page=page_index,
                title=f"page {page_index}",
            )
            if block:
                blocks.append(block)
                block_index += 1
    return blocks


def extract_pptx_v2(path: Path, source: SourceDocument) -> list[ExtractedBlock]:
    prs = Presentation(str(path))
    blocks: list[ExtractedBlock] = []
    block_index = 0
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text_parts: list[str] = []
        shape_index = 0
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = normalize_text(shape.text)
            if not text:
                continue
            shape_index += 1
            slide_text_parts.append(text)
            block = make_block(
                source=source,
                block_index=block_index,
                block_type="shape_text",
                text=text,
                slide=slide_index,
                title=f"slide {slide_index} shape {shape_index}",
            )
            if block:
                blocks.append(block)
                block_index += 1
        if slide_text_parts:
            block = make_block(
                source=source,
                block_index=block_index,
                block_type="slide",
                text="\n".join(slide_text_parts),
                slide=slide_index,
                title=f"slide {slide_index}",
            )
            if block:
                blocks.append(block)
                block_index += 1
    return blocks


def extract_html_v2(path: Path, source: SourceDocument) -> list[ExtractedBlock]:
    soup = BeautifulSoup(read_text_guess(path), "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    blocks: list[ExtractedBlock] = []
    block_index = 0
    for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th"]):
        text = normalize_text(tag.get_text(" "))
        if not text:
            continue
        heading_level = int(tag.name[1]) if tag.name.startswith("h") and tag.name[1:].isdigit() else None
        block = make_block(
            source=source,
            block_index=block_index,
            block_type="heading" if heading_level is not None else "html_text",
            text=text,
            heading_level=heading_level,
            title=text if heading_level is not None else None,
        )
        if block:
            blocks.append(block)
            block_index += 1
    if not blocks:
        text = normalize_text(soup.get_text("\n"))
        block = make_block(source=source, block_index=0, block_type="html", text=text)
        if block:
            blocks.append(block)
    return blocks


def extract_text_v2(path: Path, source: SourceDocument) -> list[ExtractedBlock]:
    text = read_text_guess(path)
    blocks: list[ExtractedBlock] = []
    block_index = 0
    parts = [part.strip() for part in text.split("\n\n") if part.strip()]
    if not parts:
        parts = [text]
    for part in parts:
        block = make_block(source=source, block_index=block_index, block_type="text", text=part)
        if block:
            blocks.append(block)
            block_index += 1
    return blocks


def extract_source(path: Path, source: SourceDocument) -> list[ExtractedBlock]:
    ext = path.suffix.lower()
    if ext == ".docx":
        return extract_docx_v2(path, source)
    if ext in {".xlsx", ".xlsb"}:
        return extract_xlsx_v2(path, source)
    if ext == ".pdf":
        return extract_pdf_v2(path, source)
    if ext == ".pptx":
        return extract_pptx_v2(path, source)
    if ext == ".html":
        return extract_html_v2(path, source)
    if ext in TEXT_EXTENSIONS:
        return extract_text_v2(path, source)
    return []


def write_report_markdown(path: Path, report: dict[str, Any]) -> None:
    lines = [
        "# Asu June Bot Extraction v2 Report",
        "",
        f"Generated: {report['generated_at']}",
        "",
        "## Summary",
        "",
    ]
    for key, value in report["summary"].items():
        lines.append(f"- {key}: {value}")
    lines.extend(["", "## By Extension", ""])
    for key, value in sorted(report["by_extension"].items()):
        lines.append(f"- {key}: {value}")
    lines.extend(["", "## By Block Type", ""])
    for key, value in sorted(report["by_block_type"].items()):
        lines.append(f"- {key}: {value}")
    lines.extend(["", "## By Document Type", ""])
    for key, value in sorted(report["by_document_type"].items()):
        lines.append(f"- {key}: {value}")
    lines.extend(["", "## Errors", ""])
    if report["errors"]:
        for item in report["errors"][:100]:
            lines.append(f"- {item['relative_path']}: {item['error']}")
    else:
        lines.append("- no errors")
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def build_report(
    *,
    sources: list[SourceDocument],
    blocks: list[ExtractedBlock],
    errors: list[dict[str, Any]],
    dry_run: bool,
    skipped_existing: int,
    total_candidates: int,
    pending_candidates: int,
) -> dict[str, Any]:
    return {
        "generated_at": utc_now(),
        "extractor_version": EXTRACTOR_VERSION,
        "dry_run": dry_run,
        "resume_supported": True,
        "summary": {
            "candidate_sources_total": total_candidates,
            "pending_sources_this_run": pending_candidates,
            "sources_extracted_this_run": len(sources),
            "sources_skipped_existing": skipped_existing,
            "blocks_extracted_this_run": len(blocks),
            "errors_this_run": len(errors),
            "docx_sources_this_run": len([src for src in sources if src.extension == ".docx"]),
            "xlsx_sources_this_run": len([src for src in sources if src.extension in {".xlsx", ".xlsb"}]),
            "pdf_sources_this_run": len([src for src in sources if src.extension == ".pdf"]),
            "blocks_with_section_this_run": len([block for block in blocks if block.section]),
            "table_row_blocks_this_run": len([block for block in blocks if block.block_type == "table_row"]),
        },
        "by_extension": dict(Counter(src.extension or "unknown" for src in sources)),
        "by_block_type": dict(Counter(block.block_type for block in blocks)),
        "by_document_type": dict(Counter(str(block.document_type or "unknown") for block in blocks)),
        "errors": errors,
    }


def write_progress(progress_path: Path, payload: dict[str, Any]) -> None:
    progress_path.parent.mkdir(parents=True, exist_ok=True)
    progress_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Asu June Bot extractor v2: independent structured extraction pipeline")
    parser.add_argument("--dry-run", action="store_true", help="Do not write output files")
    parser.add_argument("--limit", type=int, default=0, help="Limit source files")
    parser.add_argument("--path-contains", default=None, help="Process only source files whose relative path contains substring")
    parser.add_argument("--output-dir", default=DEFAULT_OUTPUT_DIR, help="Output directory")
    parser.add_argument("--reset", action="store_true", help="Delete output directory before extraction and start from scratch")
    parser.add_argument("--no-resume", action="store_true", help="Do not skip already extracted source_id records")
    parser.add_argument("--project-root", default=None, help="Override project_root without editing config.yaml")
    parser.add_argument("--source-links", default=None, help="Optional source_links.jsonl from asu_june_bot_build_source_links.py")
    parser.add_argument("--exclude-dir", action="append", default=[], help="Additional directory name to exclude. Can be repeated.")
    parser.add_argument("--exclude-path-pattern", action="append", default=[], help="Additional fnmatch path pattern to exclude. Can be repeated.")
    args = parser.parse_args()

    cfg = load_config()
    if args.project_root:
        project_root = Path(args.project_root).resolve()
        if not project_root.exists():
            raise FileNotFoundError(f"Project root not found: {project_root}")
        cfg["project_root"] = project_root.as_posix()
        cfg["project_root_path"] = project_root
    if args.source_links:
        source_links_path = resolve_work_path(cfg, args.source_links)
        cfg["source_links"] = load_source_links(source_links_path)
    if args.exclude_dir:
        cfg["exclude_dirs"] = list(cfg.get("exclude_dirs") or []) + list(args.exclude_dir)
    if args.exclude_path_pattern:
        cfg["exclude_path_patterns"] = list(cfg.get("exclude_path_patterns") or []) + list(args.exclude_path_pattern)
    output_dir = resolve_work_path(cfg, args.output_dir)
    blocks_path = output_dir / "blocks.jsonl"
    documents_path = output_dir / "documents.jsonl"
    errors_path = output_dir / "errors.jsonl"
    progress_path = output_dir / "extraction_v2_progress.json"
    report_json_path = output_dir / "extraction_v2_report.json"
    report_md_path = output_dir / "extraction_v2_report.md"

    if args.reset and not args.dry_run and output_dir.exists():
        shutil.rmtree(output_dir)

    completed_source_ids = set() if args.no_resume or args.reset else load_completed_source_ids(documents_path)

    paths = iter_source_files(cfg)
    if args.path_contains:
        needle = args.path_contains.lower()
        paths = [path for path in paths if needle in relative_to_project(cfg, path).lower()]
    if args.limit and args.limit > 0:
        paths = paths[: args.limit]

    total_candidates = len(paths)
    pending_paths: list[tuple[Path, SourceDocument]] = []
    skipped_existing = 0
    source_errors: list[dict[str, Any]] = []

    for path in paths:
        try:
            source = make_source_document(cfg, path)
            if source.source_id in completed_source_ids:
                skipped_existing += 1
                continue
            pending_paths.append((path, source))
        except Exception as exc:  # noqa: BLE001
            source_errors.append({"relative_path": relative_to_project(cfg, path), "error": repr(exc)})

    sources_this_run: list[SourceDocument] = []
    blocks_this_run: list[ExtractedBlock] = []
    errors: list[dict[str, Any]] = list(source_errors)

    for index, (path, source) in enumerate(pending_paths, start=1):
        try:
            extracted = extract_source(path, source)
            if not extracted:
                continue
            sources_this_run.append(source)
            blocks_this_run.extend(extracted)
            if not args.dry_run:
                append_jsonl(documents_path, [source.to_dict()])
                append_jsonl(blocks_path, [block.to_dict() for block in extracted])
                completed_source_ids.add(source.source_id)
                write_progress(
                    progress_path,
                    {
                        "updated_at": utc_now(),
                        "candidate_sources_total": total_candidates,
                        "pending_sources_at_start": len(pending_paths),
                        "processed_this_run": index,
                        "completed_sources_total": len(completed_source_ids),
                        "last_source": source.relative_path,
                        "last_source_id": source.source_id,
                    },
                )
        except Exception as exc:  # noqa: BLE001 - extraction should continue per file.
            error = {"relative_path": source.relative_path, "source_id": source.source_id, "error": repr(exc)}
            errors.append(error)
            if not args.dry_run:
                append_jsonl(errors_path, [error])

    report = build_report(
        sources=sources_this_run,
        blocks=blocks_this_run,
        errors=errors,
        dry_run=args.dry_run,
        skipped_existing=skipped_existing,
        total_candidates=total_candidates,
        pending_candidates=len(pending_paths),
    )

    if not args.dry_run:
        output_dir.mkdir(parents=True, exist_ok=True)
        report_json_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
        write_report_markdown(report_md_path, report)

    print(json.dumps(report, ensure_ascii=False, indent=2))
    if args.dry_run:
        print("Dry-run: no files written.")
    else:
        print(
            json.dumps(
                {
                    "documents_path": str(documents_path),
                    "blocks_path": str(blocks_path),
                    "errors_path": str(errors_path),
                    "progress_path": str(progress_path),
                    "report_json": str(report_json_path),
                    "report_md": str(report_md_path),
                },
                ensure_ascii=False,
                indent=2,
            )
        )


if __name__ == "__main__":
    main()
